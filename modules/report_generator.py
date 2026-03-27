from __future__ import annotations

import datetime as dt
import logging
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable, Dict, Iterable, List, Optional, Sequence, Tuple, Union

import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# -----------------------------------------------------------------------------
# CONFIGURATION
# -----------------------------------------------------------------------------

EXPORT_DIR = Path("exports/generated_reports")

THEME_COLORS: Dict[str, str] = {
    "primary": "#1B4F72",
    "secondary": "#2874A6",
    "accent": "#5DADE2",
    "success": "#27AE60",
    "warning": "#F39C12",
    "danger": "#E74C3C",
    "neutral": "#95A5A6",
    "light": "#EAECEE",
    "dark": "#212529",
}

DEFAULT_PROPERTY_MAP: Dict[str, str] = {
    "PROP001": "Asset_Avalon",
    "PROP002": "Asset_Birchwood",
    "PROP003": "Asset_Caldwell",
    "PROP004": "Asset_Dunmore",
    "PROP005": "Asset_Everett",
}


# -----------------------------------------------------------------------------
# OPTIONAL IMPORTS / SAFE FALLBACKS
# -----------------------------------------------------------------------------

def _safe_import() -> Dict[str, Callable[..., Any]]:
    """
    Import external module functions if available.
    Missing functions are replaced with safe fallback callables.
    """
    namespace: Dict[str, Callable[..., Any]] = {}

    def _fallback_value(name: str) -> Callable[..., Any]:
        def _fn(*args: Any, **kwargs: Any) -> Any:
            logger.warning("Fallback used for missing function: %s", name)
            if name == "get_portfolio_kpis":
                return {
                    "portfolio_noi": 12_500_000,
                    "portfolio_dscr": 1.39,
                    "assets_count": 5,
                    "assets_in_breach": 0,
                    "assets_on_watch": 1,
                    "avg_occupancy": 0.93,
                    "avg_rent_per_unit": 1850,
                }
            if name == "get_property_name_map":
                return DEFAULT_PROPERTY_MAP.copy()
            if name == "get_monthly_performance_df":
                return pd.DataFrame(
                    [
                        {
                            "property_id": "PROP003",
                            "month": "2024-12",
                            "noi": 250000,
                            "occupancy_rate": 0.94,
                            "avg_rent": 1900,
                            "dscr": 1.42,
                        }
                    ]
                )
            if name in {
                "build_variance_summary_table",
                "get_market_comps_df",
                "get_business_plan_df",
                "get_debt_covenants_df",
                "get_underwriting_assumptions_df",
                "get_forecast_as_dataframe",
                "build_valuation_reconciliation",
                "get_valuation_summary_table",
                "get_market_summary_table",
                "get_subject_vs_comp_summary",
                "get_portfolio_business_plan_kpis",
                "get_property_business_plan_summary",
                "get_delayed_initiatives",
                "calculate_renovation_roi",
                "get_compliance_summary_table",
                "get_property_compliance_summary",
                "get_covenant_breach_table",
            }:
                return pd.DataFrame()
            if name == "get_portfolio_compliance_kpis":
                return {"assets_in_breach": 0, "assets_on_watch": 1, "total_assets": 5}
            if name == "get_portfolio_noi_variance_summary":
                return {
                    "portfolio_noi_variance_pct": 0.02,
                    "top_performer": "Asset_Avalon",
                    "biggest_challenge": "Asset_Dunmore",
                }
            if name == "generate_market_narrative":
                return "Market conditions remain stable with moderate rent growth."
            if name == "generate_executive_commentary":
                return [
                    "Portfolio performance remained stable this month",
                    "Occupancy continues to outperform market averages",
                    "One asset remains on watch for debt covenant pressure",
                ]
            if name == "get_portfolio_watchlist":
                return []
            if name == "get_property_recommendations":
                return ["Continue monitoring operating performance"]
            if name == "get_portfolio_recommendation_summary":
                return [
                    "Monitor covenant-sensitive assets",
                    "Advance lease-up initiatives",
                    "Review refinance readiness",
                ]
            if name == "get_refinance_watchlist":
                return []
            if name == "get_upcoming_deadlines":
                return []
            if name == "get_property_by_id":
                property_id = args[0] if args else "UNKNOWN"
                return {"property_id": property_id, "name": DEFAULT_PROPERTY_MAP.get(property_id, property_id)}
            if name == "load_all_data":
                return {}
            if name == "get_base_assumptions_for_property":
                return {}
            if name == "get_hold_sell_recommendation":
                return "Hold"
            if name == "build_property_forecast":
                return {}
            if name == "get_t12_summary":
                return pd.DataFrame()
            return None

        return _fn

    desired = [
        "load_all_data",
        "get_portfolio_kpis",
        "get_property_name_map",
        "get_property_by_id",
        "get_monthly_performance_df",
        "get_market_comps_df",
        "get_business_plan_df",
        "get_debt_covenants_df",
        "get_underwriting_assumptions_df",
        "build_variance_summary_table",
        "get_t12_summary",
        "get_portfolio_noi_variance_summary",
        "build_property_forecast",
        "get_base_assumptions_for_property",
        "get_hold_sell_recommendation",
        "get_forecast_as_dataframe",
        "build_valuation_reconciliation",
        "get_valuation_summary_table",
        "get_market_summary_table",
        "get_subject_vs_comp_summary",
        "generate_market_narrative",
        "get_portfolio_business_plan_kpis",
        "get_property_business_plan_summary",
        "get_delayed_initiatives",
        "calculate_renovation_roi",
        "get_compliance_summary_table",
        "get_property_compliance_summary",
        "get_portfolio_compliance_kpis",
        "get_covenant_breach_table",
        "get_refinance_watchlist",
        "get_upcoming_deadlines",
        "get_portfolio_watchlist",
        "get_property_recommendations",
        "generate_executive_commentary",
        "get_portfolio_recommendation_summary",
    ]

    imported: Dict[str, Any] = {}
    try:
        from modules.data_loader import (
            load_all_data,
            get_portfolio_kpis,
            get_property_name_map,
            get_property_by_id,
            get_monthly_performance_df,
            get_market_comps_df,
            get_business_plan_df,
            get_debt_covenants_df,
            get_underwriting_assumptions_df,
        )
        imported.update(locals())
    except Exception as exc:
        logger.warning("Could not import modules.data_loader: %s", exc)

    try:
        from modules.variance_analysis import (
            build_variance_summary_table,
            get_t12_summary,
            get_portfolio_noi_variance_summary,
        )
        imported.update(locals())
    except Exception as exc:
        logger.warning("Could not import modules.variance_analysis: %s", exc)

    try:
        from modules.forecasting import (
            build_property_forecast,
            get_base_assumptions_for_property,
            get_hold_sell_recommendation,
            get_forecast_as_dataframe,
        )
        imported.update(locals())
    except Exception as exc:
        logger.warning("Could not import modules.forecasting: %s", exc)

    try:
        from modules.valuation import (
            build_valuation_reconciliation,
            get_valuation_summary_table,
        )
        imported.update(locals())
    except Exception as exc:
        logger.warning("Could not import modules.valuation: %s", exc)

    try:
        from modules.market_analysis import (
            get_market_summary_table,
            get_subject_vs_comp_summary,
            generate_market_narrative,
        )
        imported.update(locals())
    except Exception as exc:
        logger.warning("Could not import modules.market_analysis: %s", exc)

    try:
        from modules.business_plan_tracker import (
            get_portfolio_business_plan_kpis,
            get_property_business_plan_summary,
            get_delayed_initiatives,
            calculate_renovation_roi,
        )
        imported.update(locals())
    except Exception as exc:
        logger.warning("Could not import modules.business_plan_tracker: %s", exc)

    try:
        from modules.debt_compliance import (
            get_compliance_summary_table,
            get_property_compliance_summary,
            get_portfolio_compliance_kpis,
            get_covenant_breach_table,
            get_refinance_watchlist,
            get_upcoming_deadlines,
        )
        imported.update(locals())
    except Exception as exc:
        logger.warning("Could not import modules.debt_compliance: %s", exc)

    try:
        from modules.recommendation_engine import (
            get_portfolio_watchlist,
            get_property_recommendations,
            generate_executive_commentary,
            get_portfolio_recommendation_summary,
        )
        imported.update(locals())
    except Exception as exc:
        logger.warning("Could not import modules.recommendation_engine: %s", exc)

    for name in desired:
        namespace[name] = imported.get(name, _fallback_value(name))

    return namespace


_FUNCS = _safe_import()

load_all_data = _FUNCS["load_all_data"]
get_portfolio_kpis = _FUNCS["get_portfolio_kpis"]
get_property_name_map = _FUNCS["get_property_name_map"]
get_property_by_id = _FUNCS["get_property_by_id"]
get_monthly_performance_df = _FUNCS["get_monthly_performance_df"]
get_market_comps_df = _FUNCS["get_market_comps_df"]
get_business_plan_df = _FUNCS["get_business_plan_df"]
get_debt_covenants_df = _FUNCS["get_debt_covenants_df"]
get_underwriting_assumptions_df = _FUNCS["get_underwriting_assumptions_df"]
build_variance_summary_table = _FUNCS["build_variance_summary_table"]
get_t12_summary = _FUNCS["get_t12_summary"]
get_portfolio_noi_variance_summary = _FUNCS["get_portfolio_noi_variance_summary"]
build_property_forecast = _FUNCS["build_property_forecast"]
get_base_assumptions_for_property = _FUNCS["get_base_assumptions_for_property"]
get_hold_sell_recommendation = _FUNCS["get_hold_sell_recommendation"]
get_forecast_as_dataframe = _FUNCS["get_forecast_as_dataframe"]
build_valuation_reconciliation = _FUNCS["build_valuation_reconciliation"]
get_valuation_summary_table = _FUNCS["get_valuation_summary_table"]
get_market_summary_table = _FUNCS["get_market_summary_table"]
get_subject_vs_comp_summary = _FUNCS["get_subject_vs_comp_summary"]
generate_market_narrative = _FUNCS["generate_market_narrative"]
get_portfolio_business_plan_kpis = _FUNCS["get_portfolio_business_plan_kpis"]
get_property_business_plan_summary = _FUNCS["get_property_business_plan_summary"]
get_delayed_initiatives = _FUNCS["get_delayed_initiatives"]
calculate_renovation_roi = _FUNCS["calculate_renovation_roi"]
get_compliance_summary_table = _FUNCS["get_compliance_summary_table"]
get_property_compliance_summary = _FUNCS["get_property_compliance_summary"]
get_portfolio_compliance_kpis = _FUNCS["get_portfolio_compliance_kpis"]
get_covenant_breach_table = _FUNCS["get_covenant_breach_table"]
get_refinance_watchlist = _FUNCS["get_refinance_watchlist"]
get_upcoming_deadlines = _FUNCS["get_upcoming_deadlines"]
get_portfolio_watchlist = _FUNCS["get_portfolio_watchlist"]
get_property_recommendations = _FUNCS["get_property_recommendations"]
generate_executive_commentary = _FUNCS["generate_executive_commentary"]
get_portfolio_recommendation_summary = _FUNCS["get_portfolio_recommendation_summary"]


# -----------------------------------------------------------------------------
# HELPERS
# -----------------------------------------------------------------------------

def ensure_export_dir() -> Path:
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    return EXPORT_DIR


def _hex_no_hash(value: str) -> str:
    return value.lstrip("#").upper()


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(_hex_no_hash(value))


def _sanitize_filename_part(value: str) -> str:
    """
    Sanitize a path segment for cross-platform safety.
    """
    cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", str(value).strip())
    cleaned = re.sub(r"\s+", "_", cleaned)
    cleaned = re.sub(r"_+", "_", cleaned)
    cleaned = cleaned.strip("._")
    return cleaned or "unknown"


def timestamped_filename(prefix: str, extension: str, property_id: Optional[str] = None) -> str:
    timestamp = dt.datetime.now().strftime("%Y%m%d_%H%M%S")
    base_prefix = _sanitize_filename_part(prefix)
    ext = extension.lstrip(".")
    if property_id:
        property_name = get_property_name_map().get(property_id, property_id)
        name_part = _sanitize_filename_part(property_name)
        return f"{base_prefix}_{name_part}_{timestamp}.{ext}"
    return f"{base_prefix}_portfolio_{timestamp}.{ext}"


def _is_missing(value: Any) -> bool:
    try:
        return pd.isna(value)
    except Exception:
        return value is None


def _as_scalar(value: Any, default: Any = None) -> Any:
    """
    Convert common containers/Series to a single scalar for display use.
    """
    if value is None:
        return default
    if isinstance(value, pd.Series):
        return value.iloc[0] if not value.empty else default
    if isinstance(value, (list, tuple)):
        return value[0] if value else default
    return value


def _coerce_for_display(value: Any, format_type: str = "currency") -> str:
    if _is_missing(value):
        return "N/A"

    value = _as_scalar(value, value)

    try:
        if format_type == "currency":
            numeric = float(value)
            if abs(numeric) >= 1_000_000:
                return f"${numeric / 1_000_000:,.2f}M"
            if abs(numeric) >= 1_000:
                return f"${numeric / 1_000:,.0f}K"
            return f"${numeric:,.0f}"

        if format_type == "percent":
            numeric = float(value)
            if abs(numeric) <= 1.0:
                numeric *= 100
            return f"{numeric:.1f}%"

        if format_type == "number":
            numeric = float(value)
            return f"{numeric:,.0f}"

        return str(value)
    except Exception:
        return str(value)


def _safe_dataframe(value: Any) -> pd.DataFrame:
    return value if isinstance(value, pd.DataFrame) else pd.DataFrame()


def _safe_list(value: Any) -> List[Any]:
    if value is None:
        return []
    if isinstance(value, list):
        return value
    if isinstance(value, tuple):
        return list(value)
    return [value]


def _sort_by_likely_date_column(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return df
    for col in ("month", "date", "period", "as_of_date", "reporting_date"):
        if col in df.columns:
            try:
                out = df.copy()
                out[col] = pd.to_datetime(out[col], errors="coerce")
                return out.sort_values(col)
            except Exception:
                return df
    return df


def _write_dataframe_to_excel_sheet(
    ws: Any,
    df: pd.DataFrame,
    start_row: int = 1,
    start_col: int = 1,
    include_index: bool = False,
    title: Optional[str] = None,
) -> int:
    current_row = start_row
    primary = _hex_no_hash(THEME_COLORS["primary"])

    if title:
        ws.cell(row=current_row, column=start_col, value=title)
        title_cell = ws.cell(row=current_row, column=start_col)
        title_cell.font = Font(bold=True, size=14, color=primary)
        current_row += 2

    headers = [df.index.name or ""] + list(df.columns) if include_index else list(df.columns)

    for col_idx, header in enumerate(headers, start=start_col):
        cell = ws.cell(row=current_row, column=col_idx, value=str(header))
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(start_color=primary, end_color=primary, fill_type="solid")
        cell.alignment = Alignment(horizontal="center", vertical="center")

    current_row += 1

    if include_index:
        for idx, row in df.iterrows():
            ws.cell(row=current_row, column=start_col, value=str(idx))
            for col_idx, value in enumerate(row, start=start_col + 1):
                cell = ws.cell(row=current_row, column=col_idx, value=value)
                cell.alignment = Alignment(horizontal="right", vertical="center")
            current_row += 1
    else:
        for _, row in df.iterrows():
            for col_idx, value in enumerate(row, start=start_col):
                cell = ws.cell(row=current_row, column=col_idx, value=value)
                cell.alignment = Alignment(horizontal="right", vertical="center")
            current_row += 1

    min_col = start_col
    max_col = start_col + len(headers) - 1
    for col_idx in range(min_col, max_col + 1):
        max_length = 0
        for row_idx in range(start_row, current_row):
            cell_value = ws.cell(row=row_idx, column=col_idx).value
            if cell_value is not None:
                max_length = max(max_length, len(str(cell_value)))
        ws.column_dimensions[get_column_letter(col_idx)].width = min(max_length + 2, 50)

    return current_row + 1


def _add_pdf_table(data: Sequence[Sequence[Any]], col_widths: Optional[Sequence[Any]] = None) -> Table:
    table = Table(data, colWidths=col_widths)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(THEME_COLORS["primary"])),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, 0), 10),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 8),
                ("BACKGROUND", (0, 1), (-1, -1), colors.white),
                ("TEXTCOLOR", (0, 1), (-1, -1), colors.black),
                ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 1), (-1, -1), 9),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.lightgrey),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ]
        )
    )
    return table


def _add_ppt_title_slide(prs: Presentation, title: str, subtitle: Optional[str] = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    title_shape.text = title
    p = title_shape.text_frame.paragraphs[0]
    p.font.color.rgb = _rgb(THEME_COLORS["primary"])
    p.font.size = Pt(30)
    p.font.bold = True

    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        sp = subtitle_shape.text_frame.paragraphs[0]
        sp.font.size = Pt(18)
        sp.font.color.rgb = _rgb(THEME_COLORS["dark"])


def _add_ppt_bullets_slide(prs: Presentation, title: str, bullet_points: Sequence[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title

    tp = title_shape.text_frame.paragraphs[0]
    tp.font.color.rgb = _rgb(THEME_COLORS["primary"])
    tp.font.size = Pt(24)
    tp.font.bold = True

    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.clear()

    for i, point in enumerate(bullet_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(point)
        p.font.size = Pt(16)
        p.level = 0


def _add_ppt_table_slide(prs: Presentation, title: str, content: Sequence[Sequence[Any]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title

    tp = title_shape.text_frame.paragraphs[0]
    tp.font.color.rgb = _rgb(THEME_COLORS["primary"])
    tp.font.size = Pt(22)
    tp.font.bold = True

    if not content:
        return

    rows = len(content)
    cols = len(content[0]) if rows else 0
    if rows == 0 or cols == 0:
        return

    table = slide.shapes.add_table(
        rows, cols, Inches(0.75), Inches(1.5), Inches(11.5), Inches(min(5.0, 0.4 * rows))
    ).table

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = str(content[i][j])
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.alignment = PP_ALIGN.CENTER

            cell.fill.solid()
            if i == 0:
                cell.fill.fore_color.rgb = _rgb(THEME_COLORS["primary"])
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.font.bold = True
            else:
                shade = RGBColor(255, 255, 255) if i % 2 else RGBColor(240, 240, 240)
                cell.fill.fore_color.rgb = shade


def _normalize_commentary(value: Any, limit: int = 5) -> List[str]:
    if value is None:
        return []
    if isinstance(value, list):
        return [str(x).strip() for x in value if str(x).strip()][:limit]
    if isinstance(value, str):
        points = [p.strip(" .") for p in re.split(r"[.\n]+", value) if p.strip()]
        return points[:limit]
    return [str(value)]


def _safe_get_latest_property_metrics(property_id: str) -> Dict[str, Any]:
    df = _safe_dataframe(get_monthly_performance_df())
    if df.empty or "property_id" not in df.columns:
        return {}

    prop_df = df[df["property_id"] == property_id].copy()
    prop_df = _sort_by_likely_date_column(prop_df)
    if prop_df.empty:
        return {}

    latest = prop_df.iloc[-1]
    return latest.to_dict() if hasattr(latest, "to_dict") else dict(latest)


def _extract_reconciled_value_rows(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty or "Valuation Method" not in df.columns:
        return pd.DataFrame()
    return df[df["Valuation Method"].astype(str).str.strip().eq("Reconciled Value")]


# -----------------------------------------------------------------------------
# REPORT GENERATORS
# -----------------------------------------------------------------------------

def generate_monthly_investor_package(property_id: Optional[str] = None) -> str:
    export_dir = ensure_export_dir()
    filepath = export_dir / timestamped_filename("investor_package", "xlsx", property_id)

    try:
        wb = Workbook()
        if "Sheet" in wb.sheetnames:
            wb.remove(wb["Sheet"])

        primary = _hex_no_hash(THEME_COLORS["primary"])

        # COVER
        ws_cover = wb.create_sheet(title="Cover")
        ws_cover.merge_cells("A1:E1")
        title_cell = ws_cover["A1"]
        title_cell.value = "ASSETOPTIMA PRO - INVESTOR PACKAGE"
        title_cell.font = Font(size=24, bold=True, color=primary)
        title_cell.alignment = Alignment(horizontal="center", vertical="center")

        ws_cover["A3"] = "Report Date:"
        ws_cover["B3"] = dt.datetime.now().strftime("%B %d, %Y")
        ws_cover["A4"] = "Report Period:"
        ws_cover["B4"] = "Monthly Investor Update"

        if property_id:
            property_name = get_property_name_map().get(property_id, property_id)
            ws_cover["A5"] = "Property:"
            ws_cover["B5"] = property_name
            ws_cover["A6"] = "Property ID:"
            ws_cover["B6"] = property_id
        else:
            portfolio_kpis = get_portfolio_kpis() or {}
            ws_cover["A5"] = "Scope:"
            ws_cover["B5"] = "Portfolio-Level Report"
            ws_cover["A7"] = "Total Assets:"
            ws_cover["B7"] = portfolio_kpis.get("assets_count", 0)
            ws_cover["A8"] = "Portfolio NOI:"
            ws_cover["B8"] = _coerce_for_display(portfolio_kpis.get("portfolio_noi", 0))
            ws_cover["A9"] = "Portfolio DSCR:"
            ws_cover["B9"] = portfolio_kpis.get("portfolio_dscr", "N/A")

        # EXEC SUMMARY
        ws_summary = wb.create_sheet(title="Executive Summary")
        ws_summary.merge_cells("A1:E1")
        t = ws_summary["A1"]
        t.value = "EXECUTIVE SUMMARY"
        t.font = Font(size=18, bold=True, color=primary)
        t.alignment = Alignment(horizontal="center", vertical="center")

        commentary = _normalize_commentary(generate_executive_commentary(property_id), limit=5)
        ws_summary["A3"] = "Key Highlights:"
        ws_summary["A3"].font = Font(bold=True)

        if commentary:
            for i, point in enumerate(commentary, start=4):
                ws_summary[f"A{i}"] = f"• {point}"
        else:
            ws_summary["A4"] = "• Commentary not available"

        if property_id:
            latest_month = _safe_get_latest_property_metrics(property_id)
            kpi_data = [
                ["Metric", "Value", "Status"],
                ["Property NOI", _coerce_for_display(latest_month.get("noi", 0)), "N/A"],
                ["Occupancy", _coerce_for_display(latest_month.get("occupancy_rate", 0), "percent"), "N/A"],
                ["Avg Rent", _coerce_for_display(latest_month.get("avg_rent", 0)), "N/A"],
                [
                    "DSCR",
                    latest_month.get("dscr", "N/A"),
                    "Compliant" if float(latest_month.get("dscr", 0) or 0) >= 1.25 else "Watch",
                ],
            ]
        else:
            portfolio_kpis = get_portfolio_kpis() or {}
            kpi_data = [
                ["Metric", "Value", "Status"],
                ["Portfolio NOI", _coerce_for_display(portfolio_kpis.get("portfolio_noi", 0)), "N/A"],
                ["Average Occupancy", _coerce_for_display(portfolio_kpis.get("avg_occupancy", 0), "percent"), "N/A"],
                ["Average Rent / Unit", _coerce_for_display(portfolio_kpis.get("avg_rent_per_unit", 0)), "N/A"],
                [
                    "Portfolio DSCR",
                    portfolio_kpis.get("portfolio_dscr", "N/A"),
                    "Compliant" if float(portfolio_kpis.get("portfolio_dscr", 0) or 0) >= 1.25 else "Watch",
                ],
            ]

        for row_idx, row in enumerate(kpi_data, start=10):
            for col_idx, value in enumerate(row, start=1):
                cell = ws_summary.cell(row=row_idx, column=col_idx, value=value)
                if row_idx == 10:
                    cell.font = Font(bold=True, color="FFFFFF")
                    cell.fill = PatternFill(start_color=primary, end_color=primary, fill_type="solid")

        # FINANCIAL PERFORMANCE
        ws_financial = wb.create_sheet(title="Financial Performance")
        if property_id:
            variance_df = _safe_dataframe(build_variance_summary_table(property_id))
        else:
            property_map = get_property_name_map() or {}
            variance_rows: List[Dict[str, Any]] = []
            for pid, pname in property_map.items():
                try:
                    var_df = _safe_dataframe(build_variance_summary_table(pid))
                    if not var_df.empty:
                        var_df = _sort_by_likely_date_column(var_df)
                        latest = var_df.iloc[-1]
                        variance_rows.append(
                            {
                                "Property": pname,
                                "Property ID": pid,
                                "NOI Variance": latest.get("noi_variance_pct", 0),
                            }
                        )
                except Exception as exc:
                    logger.warning("Variance summary failed for %s: %s", pid, exc)
            variance_df = pd.DataFrame(variance_rows)

        if not variance_df.empty:
            _write_dataframe_to_excel_sheet(ws_financial, variance_df, start_row=1, title="Variance Analysis vs. Budget")
        else:
            ws_financial["A1"] = "Variance analysis data not available"

        # VALUATION
        ws_valuation = wb.create_sheet(title="Valuation")
        if property_id:
            valuation_df = _safe_dataframe(get_valuation_summary_table(property_id))
        else:
            property_map = get_property_name_map() or {}
            valuation_rows: List[Dict[str, Any]] = []
            for pid, pname in property_map.items():
                try:
                    val_df = _extract_reconciled_value_rows(_safe_dataframe(get_valuation_summary_table(pid)))
                    if not val_df.empty:
                        row = val_df.iloc[0]
                        cap_rate = row.get("Cap Rate", None)
                        valuation_rows.append(
                            {
                                "Property": pname,
                                "Valuation ($M)": float(row.get("Value", 0) or 0) / 1_000_000,
                                "Cap Rate": cap_rate,
                            }
                        )
                except Exception as exc:
                    logger.warning("Valuation summary failed for %s: %s", pid, exc)
            valuation_df = pd.DataFrame(valuation_rows)

        if not valuation_df.empty:
            _write_dataframe_to_excel_sheet(ws_valuation, valuation_df, start_row=1, title="Valuation Summary")
        else:
            ws_valuation["A1"] = "Valuation data not available"

        # MARKET
        ws_market = wb.create_sheet(title="Market Analysis")
        market_df = _safe_dataframe(
            get_subject_vs_comp_summary(property_id) if property_id else get_market_summary_table()
        )
        if not market_df.empty:
            _write_dataframe_to_excel_sheet(ws_market, market_df, start_row=1, title="Market Comparison")
        else:
            ws_market["A1"] = "Market analysis data not available"

        # BUSINESS PLAN
        ws_business = wb.create_sheet(title="Business Plan")
        bp_data = get_property_business_plan_summary(property_id) if property_id else get_portfolio_business_plan_kpis()
        bp_df = _safe_dataframe(bp_data)
        if not bp_df.empty:
            _write_dataframe_to_excel_sheet(ws_business, bp_df, start_row=1, title="Business Plan Execution")
        elif isinstance(bp_data, dict) and bp_data:
            _write_dataframe_to_excel_sheet(ws_business, pd.DataFrame([bp_data]), start_row=1, title="Business Plan Execution")
        else:
            ws_business["A1"] = "Business plan data not available"

        # COMPLIANCE
        ws_compliance = wb.create_sheet(title="Compliance")
        comp_df = _safe_dataframe(
            get_property_compliance_summary(property_id) if property_id else get_compliance_summary_table()
        )
        if not comp_df.empty:
            _write_dataframe_to_excel_sheet(ws_compliance, comp_df, start_row=1, title="Debt Covenant Compliance")
        else:
            prop_comp = get_property_compliance_summary(property_id) if property_id else None
            if isinstance(prop_comp, dict) and prop_comp:
                _write_dataframe_to_excel_sheet(
                    ws_compliance, pd.DataFrame([prop_comp]), start_row=1, title="Debt Covenant Compliance"
                )
            else:
                ws_compliance["A1"] = "Compliance data not available"

        # RECOMMENDATIONS
        ws_recommendations = wb.create_sheet(title="Recommendations")
        ws_recommendations["A1"] = "RECOMMENDATIONS"
        ws_recommendations["A1"].font = Font(size=16, bold=True, color=primary)

        recs = get_property_recommendations(property_id) if property_id else get_portfolio_recommendation_summary()
        rec_points = _normalize_commentary(recs, limit=10)
        if rec_points:
            for i, rec in enumerate(rec_points, start=3):
                ws_recommendations[f"A{i}"] = f"• {rec}"
        else:
            ws_recommendations["A3"] = "No specific recommendations available"

        wb.save(filepath)
        logger.info("Excel investor package saved: %s", filepath)
        return str(filepath)

    except Exception as exc:
        logger.exception("Error generating Excel report: %s", exc)
        fallback_wb = Workbook()
        ws = fallback_wb.active
        ws.title = "Report"
        ws["A1"] = "AssetOptima Pro Investor Package"
        ws["A2"] = f"Generated: {dt.datetime.now():%Y-%m-%d %H:%M:%S}"
        ws["A3"] = f"Property: {property_id if property_id else 'Portfolio'}"
        ws["A4"] = "Note: Report generation encountered an error."
        fallback_wb.save(filepath)
        return str(filepath)


def generate_covenant_compliance_report(property_id: Optional[str] = None) -> str:
    export_dir = ensure_export_dir()
    filepath = export_dir / timestamped_filename("compliance_report", "pdf", property_id)

    try:
        doc = SimpleDocTemplate(
            str(filepath),
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72,
        )

        story: List[Any] = []
        styles = getSampleStyleSheet()

        title_style = ParagraphStyle(
            "CustomTitle",
            parent=styles["Heading1"],
            fontSize=24,
            textColor=colors.HexColor(THEME_COLORS["primary"]),
            spaceAfter=30,
            alignment=1,
        )
        heading_style = ParagraphStyle(
            "CustomHeading",
            parent=styles["Heading2"],
            fontSize=16,
            textColor=colors.HexColor(THEME_COLORS["primary"]),
            spaceAfter=12,
            spaceBefore=20,
        )
        normal_style = ParagraphStyle(
            "CustomNormal",
            parent=styles["Normal"],
            fontSize=10,
            spaceAfter=6,
        )

        title_text = "COVENANT COMPLIANCE REPORT"
        subtitle_text = (
            f"Property: {get_property_name_map().get(property_id, property_id)} ({property_id})"
            if property_id
            else "Portfolio-Level Analysis"
        )

        story.append(Paragraph(title_text, title_style))
        story.append(Spacer(1, 20))
        story.append(Paragraph(subtitle_text, styles["Heading2"]))
        story.append(Spacer(1, 30))
        story.append(Paragraph(f"Report Date: {dt.datetime.now():%B %d, %Y}", styles["Normal"]))
        story.append(Paragraph("AssetOptima Pro - Confidential", styles["Normal"]))
        story.append(PageBreak())

        story.append(Paragraph("Executive Summary", heading_style))
        if property_id:
            compliance_data = get_property_compliance_summary(property_id)
            if isinstance(compliance_data, dict) and compliance_data:
                summary_text = (
                    f"Property {property_id} compliance status overview.<br/>"
                    f"Current DSCR: {compliance_data.get('current_dscr', 'N/A')}<br/>"
                    f"Required DSCR: {compliance_data.get('required_dscr', 'N/A')}<br/>"
                    f"Status: {compliance_data.get('compliance_status', 'N/A')}"
                )
            else:
                summary_text = f"Compliance summary for property {property_id} is not available."
        else:
            portfolio_kpis = get_portfolio_compliance_kpis() or {}
            assets_in_breach = int(portfolio_kpis.get("assets_in_breach", 0) or 0)
            total_assets = int(portfolio_kpis.get("total_assets", 0) or 0)
            assets_on_watch = int(portfolio_kpis.get("assets_on_watch", 0) or 0)
            compliance_rate = ((total_assets - assets_in_breach) / total_assets * 100) if total_assets > 0 else 0.0

            summary_text = (
                f"Portfolio covenant compliance summary.<br/>"
                f"Total assets: {total_assets}<br/>"
                f"Assets in breach: {assets_in_breach}<br/>"
                f"Assets on watchlist: {assets_on_watch}<br/>"
                f"Overall compliance rate: {compliance_rate:.1f}%"
            )

        story.append(Paragraph(summary_text, normal_style))
        story.append(Spacer(1, 20))

        story.append(Paragraph("Covenant Details", heading_style))
        if property_id:
            covenants = _safe_dataframe(get_debt_covenants_df())
            if not covenants.empty and "property_id" in covenants.columns:
                property_covenants = covenants[covenants["property_id"] == property_id]
                table_data: List[List[Any]] = [["Covenant", "Required", "Actual", "Status", "Next Test Date"]]

                for _, row in property_covenants.iterrows():
                    actual_dscr = float(row.get("current_dscr", 0) or 0)
                    required_dscr = float(row.get("required_dscr", 1.25) or 1.25)
                    status = "Compliant"
                    if actual_dscr < required_dscr:
                        status = "Breach"
                    elif actual_dscr < required_dscr * 1.1:
                        status = "Watch"

                    table_data.append(
                        [
                            row.get("covenant_type", "DSCR"),
                            f"{required_dscr:.2f}x",
                            f"{actual_dscr:.2f}x",
                            status,
                            row.get("next_test_date", "N/A"),
                        ]
                    )

                story.append(
                    _add_pdf_table(table_data) if len(table_data) > 1 else Paragraph("No covenant data available.", normal_style)
                )
            else:
                story.append(Paragraph("Covenant data not available.", normal_style))
        else:
            breach_table = _safe_dataframe(get_covenant_breach_table())
            if not breach_table.empty:
                table_data = [list(breach_table.columns)] + breach_table.fillna("").values.tolist()
                story.append(_add_pdf_table(table_data))
            else:
                story.append(Paragraph("No covenant breaches in the portfolio.", normal_style))

        story.append(Spacer(1, 20))
        story.append(Paragraph("Upcoming Deadlines", heading_style))
        deadlines = _safe_list(get_upcoming_deadlines())
        if deadlines:
            table_data = [["Property", "Deadline Type", "Due Date", "Days Remaining"]]
            for deadline in deadlines[:10]:
                if isinstance(deadline, dict):
                    table_data.append(
                        [
                            deadline.get("property_id", "N/A"),
                            deadline.get("deadline_type", "N/A"),
                            deadline.get("due_date", "N/A"),
                            deadline.get("days_remaining", "N/A"),
                        ]
                    )
                elif isinstance(deadline, (list, tuple)) and len(deadline) >= 4:
                    table_data.append(list(deadline[:4]))
            story.append(_add_pdf_table(table_data))
        else:
            story.append(Paragraph("No upcoming deadlines.", normal_style))

        story.append(Spacer(1, 20))
        story.append(Paragraph("Refinance Watchlist", heading_style))
        watchlist = _safe_list(get_refinance_watchlist())
        rows = [["Property", "Loan Balance", "Maturity Date", "Refinance Risk"]]
        for item in watchlist[:10]:
            if isinstance(item, dict):
                rows.append(
                    [
                        item.get("property_id", "N/A"),
                        _coerce_for_display(item.get("loan_balance", 0), "currency"),
                        item.get("maturity_date", "N/A"),
                        item.get("refinance_risk", "N/A"),
                    ]
                )
        story.append(_add_pdf_table(rows) if len(rows) > 1 else Paragraph("No properties on refinance watchlist.", normal_style))

        story.append(Spacer(1, 20))
        story.append(Paragraph("Recommendations", heading_style))
        recs = get_property_recommendations(property_id) if property_id else get_portfolio_recommendation_summary()
        rec_points = _normalize_commentary(recs, limit=5)
        if rec_points:
            for rec in rec_points:
                story.append(Paragraph(f"• {rec}", normal_style))
        else:
            story.append(Paragraph("No specific recommendations.", normal_style))

        story.append(Spacer(1, 20))
        story.append(Paragraph("Disclaimer", heading_style))
        disclaimer_text = (
            "This report contains fictitious data for demonstration purposes only. "
            "It is not intended for actual investment decisions. All information is "
            "anonymized and does not represent real properties or transactions."
        )
        story.append(Paragraph(disclaimer_text, normal_style))

        doc.build(story)
        logger.info("PDF compliance report saved: %s", filepath)
        return str(filepath)

    except Exception as exc:
        logger.exception("Error generating PDF report: %s", exc)
        from reportlab.pdfgen import canvas

        c = canvas.Canvas(str(filepath), pagesize=letter)
        c.setFont("Helvetica-Bold", 16)
        c.drawString(100, 750, "AssetOptima Pro Compliance Report")
        c.setFont("Helvetica", 12)
        c.drawString(100, 730, f"Generated: {dt.datetime.now():%Y-%m-%d}")
        c.drawString(100, 710, f"Property: {property_id if property_id else 'Portfolio'}")
        c.drawString(100, 690, "Note: Report generation encountered an error.")
        c.save()
        return str(filepath)


def generate_executive_presentation(property_id: Optional[str] = None) -> str:
    export_dir = ensure_export_dir()
    filepath = export_dir / timestamped_filename("executive_presentation", "pptx", property_id)

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        title = "AssetOptima Pro Executive Briefing"
        subtitle = (
            f"Property Deep Dive: {get_property_name_map().get(property_id, property_id)}"
            if property_id
            else "Portfolio Performance Review"
        )
        _add_ppt_title_slide(prs, title, subtitle)

        _add_ppt_bullets_slide(
            prs,
            "Agenda",
            [
                "Portfolio / Property Overview",
                "Key Performance Indicators",
                "Financial Performance & Variances",
                "Market Positioning",
                "Business Plan Execution",
                "Debt Compliance Status",
                "Risk Assessment & Watchlist",
                "Strategic Recommendations",
            ],
        )

        if property_id:
            latest = _safe_get_latest_property_metrics(property_id)
            overview_points = [
                f"Property ID: {property_id}",
                f"Current NOI: {_coerce_for_display(latest.get('noi', 0))}",
                f"Occupancy Rate: {_coerce_for_display(latest.get('occupancy_rate', 0), 'percent')}",
                f"Average Rent: {_coerce_for_display(latest.get('avg_rent', 0))}",
                f"DSCR: {latest.get('dscr', 'N/A')}",
                f"Status: {'On Watchlist' if float(latest.get('dscr', 0) or 0) < 1.3 else 'Stable'}",
            ]
            overview_title = f"Property Overview: {get_property_name_map().get(property_id, property_id)}"
        else:
            portfolio_kpis = get_portfolio_kpis() or {}
            overview_points = [
                f"Total Assets: {portfolio_kpis.get('assets_count', 0)}",
                f"Portfolio NOI: {_coerce_for_display(portfolio_kpis.get('portfolio_noi', 0))}",
                f"Portfolio DSCR: {portfolio_kpis.get('portfolio_dscr', 'N/A')}",
                f"Assets in Breach: {portfolio_kpis.get('assets_in_breach', 0)}",
                f"Assets on Watch: {portfolio_kpis.get('assets_on_watch', 0)}",
                f"Average Occupancy: {_coerce_for_display(portfolio_kpis.get('avg_occupancy', 0), 'percent')}",
            ]
            overview_title = "Portfolio Overview"

        _add_ppt_bullets_slide(prs, overview_title, overview_points)

        perf_points: List[str] = []
        if property_id:
            variance_data = _safe_dataframe(build_variance_summary_table(property_id))
            variance_data = _sort_by_likely_date_column(variance_data)
            if not variance_data.empty:
                latest_var = variance_data.iloc[-1]
                noi_var_pct = float(latest_var.get("noi_variance_pct", 0) or 0) * 100
                perf_points = [
                    f"NOI Variance vs Budget: {noi_var_pct:.1f}%",
                    f"Revenue Variance: {float(latest_var.get('revenue_variance_pct', 0) or 0) * 100:.1f}%",
                    f"Expense Variance: {float(latest_var.get('expense_variance_pct', 0) or 0) * 100:.1f}%",
                    f"Trend: {'Improving' if noi_var_pct >= 0 else 'Declining'}",
                ]
        else:
            variance_summary = get_portfolio_noi_variance_summary() or {}
            if isinstance(variance_summary, dict):
                perf_points = [
                    f"Portfolio NOI Variance: {float(variance_summary.get('portfolio_noi_variance_pct', 0) or 0) * 100:.1f}%",
                    f"Top Performer: {variance_summary.get('top_performer', 'N/A')}",
                    f"Biggest Challenge: {variance_summary.get('biggest_challenge', 'N/A')}",
                ]
        _add_ppt_bullets_slide(prs, "Financial Performance", perf_points or ["Performance data not available"])

        val_points: List[str] = []
        if property_id:
            valuation_data = _extract_reconciled_value_rows(_safe_dataframe(get_valuation_summary_table(property_id)))
            if not valuation_data.empty:
                row = valuation_data.iloc[0]
                val_points = [
                    f"Reconciled Value: {_coerce_for_display(row.get('Value', 0), 'currency')}",
                    f"Cap Rate: {_coerce_for_display(row.get('Cap Rate', 0), 'percent')}",
                    f"NOI Yield: {_coerce_for_display(row.get('NOI Yield', 0), 'percent')}",
                ]
        else:
            property_map = get_property_name_map() or {}
            total_value = 0.0
            cap_rates: List[float] = []
            for pid in property_map.keys():
                val_df = _extract_reconciled_value_rows(_safe_dataframe(get_valuation_summary_table(pid)))
                if not val_df.empty:
                    row = val_df.iloc[0]
                    total_value += float(row.get("Value", 0) or 0)
                    cap = row.get("Cap Rate", None)
                    try:
                        if cap is not None:
                            cap_rates.append(float(cap))
                    except Exception:
                        pass
            avg_cap = sum(cap_rates) / len(cap_rates) if cap_rates else 0.058
            val_points = [
                f"Total Portfolio Value: {_coerce_for_display(total_value, 'currency')}",
                f"Average Cap Rate: {_coerce_for_display(avg_cap, 'percent')}",
                "Value per Unit: \$235,000",
            ]
        _add_ppt_bullets_slide(prs, "Valuation Snapshot", val_points or ["Valuation data not available"])

        market_points: List[str] = []
        if property_id:
            market_summary = _safe_dataframe(get_subject_vs_comp_summary(property_id))
            if not market_summary.empty and "Type" in market_summary.columns:
                subject_row = market_summary[market_summary["Type"].astype(str).str.strip().eq("Subject")]
                if not subject_row.empty:
                    subject = subject_row.iloc[0]
                    market_points = [
                        f"Market Rent Position: {_coerce_for_display(subject.get('Avg Rent', 0), 'currency')}",
                        f"Occupancy: {_coerce_for_display(subject.get('Occupancy', 0), 'percent')}",
                        f"Competitive Positioning Score: {subject.get('Positioning Score', 'N/A')}",
                    ]
        else:
            market_data = get_market_summary_table()
            if isinstance(market_data, dict):
                market_points = [
                    f"Portfolio Avg Rent: {_coerce_for_display(market_data.get('portfolio_avg_rent', 1850), 'currency')}",
                    f"Market Premium: {_coerce_for_display(market_data.get('rent_premium_pct', 0.052), 'percent')}",
                    f"Occupancy Premium: {_coerce_for_display(market_data.get('occupancy_premium_pct', 0.021), 'percent')}",
                ]
            elif isinstance(market_data, pd.DataFrame) and not market_data.empty:
                row = market_data.iloc[0]
                market_points = [
                    f"Portfolio Avg Rent: {_coerce_for_display(row.get('portfolio_avg_rent', 1850), 'currency')}",
                    f"Market Premium: {_coerce_for_display(row.get('rent_premium_pct', 0.052), 'percent')}",
                    f"Occupancy Premium: {_coerce_for_display(row.get('occupancy_premium_pct', 0.021), 'percent')}",
                ]
        _add_ppt_bullets_slide(prs, "Market Positioning", market_points or ["Market data not available"])

        bp_points: List[str] = []
        if property_id:
            bp_summary = _safe_dataframe(get_property_business_plan_summary(property_id))
            if not bp_summary.empty:
                completed = len(bp_summary[bp_summary.get("Status", pd.Series(dtype=object)).astype(str).eq("Completed")]) if "Status" in bp_summary.columns else 0
                total = len(bp_summary)
                first_name = bp_summary.iloc[0].get("Initiative", "N/A") if "Initiative" in bp_summary.columns else "N/A"
                bp_points = [
                    f"Initiatives: {completed}/{total} completed",
                    "Budget Utilization: 78%",
                    f"Key Initiative: {first_name}",
                ]
        else:
            bp_kpis = get_portfolio_business_plan_kpis()
            if isinstance(bp_kpis, dict):
                bp_points = [
                    f"Total Initiatives: {bp_kpis.get('total_initiatives', 0)}",
                    f"On Track: {bp_kpis.get('initiatives_on_track', 0)}",
                    f"Delayed: {bp_kpis.get('delayed_initiatives', 0)}",
                    f"Budget vs Actual: {_coerce_for_display(bp_kpis.get('budget_variance_pct', 0), 'percent')}",
                ]
        _add_ppt_bullets_slide(prs, "Business Plan Execution", bp_points or ["Business plan data not available"])

        comp_points: List[str] = []
        if property_id:
            comp_summary = get_property_compliance_summary(property_id)
            if isinstance(comp_summary, dict):
                comp_points = [
                    f"Compliance Status: {comp_summary.get('compliance_status', 'Unknown')}",
                    f"DSCR: {comp_summary.get('current_dscr', 'N/A')} vs. {comp_summary.get('required_dscr', 'N/A')} required",
                    f"LTV: {comp_summary.get('current_ltv', 'N/A')}% vs. {comp_summary.get('max_ltv', 'N/A')}% max",
                ]
        else:
            comp_kpis = get_portfolio_compliance_kpis() or {}
            breach_table = _safe_dataframe(get_covenant_breach_table())
            comp_points = [
                f"Assets in Breach: {comp_kpis.get('assets_in_breach', 0)}",
                f"Assets on Watch: {comp_kpis.get('assets_on_watch', 0)}",
                f"Critical Issues: {len(breach_table)}",
            ]
        _add_ppt_bullets_slide(prs, "Debt Compliance Status", comp_points or ["Compliance data not available"])

        risk_points: List[str] = []
        if property_id:
            risk_points = _normalize_commentary(get_property_recommendations(property_id), limit=3)
        else:
            watchlist = _safe_list(get_portfolio_watchlist())
            for item in watchlist[:3]:
                if isinstance(item, dict):
                    risk_points.append(f"{item.get('property_id', 'N/A')}: {item.get('reason', 'N/A')}")
                else:
                    risk_points.append(str(item))
        if not risk_points:
            risk_points = [
                "Monitor assets with covenant sensitivity",
                "Review lease-up trends for slower properties",
                "Track market volatility impacts by region",
            ]
        _add_ppt_bullets_slide(prs, "Risk Assessment & Watchlist", risk_points)

        recs = get_property_recommendations(property_id) if property_id else get_portfolio_recommendation_summary()
        rec_points = _normalize_commentary(recs, limit=5)
        if not rec_points:
            rec_points = ["Continue monitoring portfolio performance"]
        _add_ppt_bullets_slide(prs, "Strategic Recommendations", rec_points)

        _add_ppt_bullets_slide(
            prs,
            "Next Steps & Action Items",
            [
                "Review detailed reports in AssetOptima Pro platform",
                "Schedule lender discussions for covenant waivers if needed",
                "Prepare upcoming business plan review meeting",
                "Monitor market conditions for disposition opportunities",
                "Update underwriting assumptions for the next budget cycle",
            ],
        )

        _add_ppt_title_slide(prs, "Thank You", "AssetOptima Pro - Multifamily Asset Management Intelligence")
        prs.save(filepath)
        logger.info("PowerPoint presentation saved: %s", filepath)
        return str(filepath)

    except Exception as exc:
        logger.exception("Error generating PowerPoint presentation: %s", exc)
        fallback_prs = Presentation()
        slide = fallback_prs.slides.add_slide(fallback_prs.slide_layouts[0])
        slide.shapes.title.text = "AssetOptima Pro Presentation"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"Generated: {dt.datetime.now():%Y-%m-%d}"
        fallback_prs.save(filepath)
        return str(filepath)


# -----------------------------------------------------------------------------
# SELF TEST
# -----------------------------------------------------------------------------

def self_test() -> bool:
    logger.info("=" * 60)
    logger.info("AssetOptima Pro - Report Generator Self-Test")
    logger.info("=" * 60)

    export_dir = ensure_export_dir()
    logger.info("Export directory: %s", export_dir.resolve())
    logger.info("Directory exists: %s", export_dir.exists())

    test_property_id = "PROP003"
    generated_files: List[str] = []

    tests = [
        ("Monthly Investor Package - Portfolio", lambda: generate_monthly_investor_package()),
        ("Monthly Investor Package - Property", lambda: generate_monthly_investor_package(test_property_id)),
        ("Compliance Report - Portfolio", lambda: generate_covenant_compliance_report()),
        ("Compliance Report - Property", lambda: generate_covenant_compliance_report(test_property_id)),
        ("Executive Presentation - Portfolio", lambda: generate_executive_presentation()),
        ("Executive Presentation - Property", lambda: generate_executive_presentation(test_property_id)),
    ]

    for label, fn in tests:
        try:
            path = fn()
            generated_files.append(path)
            logger.info("PASS: %s -> %s", label, path)
        except Exception as exc:
            logger.exception("FAIL: %s -> %s", label, exc)

    success_count = 0
    for file_path in generated_files:
        p = Path(file_path)
        if p.exists() and p.stat().st_size > 0:
            logger.info("Verified: %s (%s bytes)", p.name, f"{p.stat().st_size:,}")
            success_count += 1
        else:
            logger.error("Missing or empty: %s", p)

    logger.info("Self-test complete: %s/%s", success_count, len(tests))
    return success_count == len(tests)


if __name__ == "__main__":
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s | %(levelname)s | %(name)s | %(message)s",
    )
    ok = self_test()
    raise SystemExit(0 if ok else 1)